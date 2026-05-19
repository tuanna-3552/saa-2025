#!/usr/bin/env python3
"""
extract-pptx.py — Convert a .pptx file to structured text for slide re-generation.

Usage:
    python3 extract-pptx.py <input.pptx> [output.md]

Output:
    Markdown file with per-slide content blocks that can be fed to generate-slide skill.

Requires:
    pip install python-pptx
"""

import sys
import os

def check_dependency():
    try:
        from pptx import Presentation  # noqa
    except ImportError:
        print("ERROR: python-pptx not installed.")
        print("Install with: pip install python-pptx")
        print("Or if using .venv: ~/.claude/skills/.venv/bin/pip install python-pptx")
        sys.exit(1)

def extract(pptx_path: str, output_path: str | None = None) -> str:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation(pptx_path)
    lines = [f"# Extracted from: {os.path.basename(pptx_path)}\n"]

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n---\n\n## Slide {i}\n")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            for para in text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Detect heading by font size or placeholder type
                is_heading = False
                if shape.shape_type == 13:  # Placeholder
                    ph = shape.placeholder_format
                    if ph and ph.idx in (0, 1):  # title / body title
                        is_heading = True
                elif para.runs and para.runs[0].font.size and para.runs[0].font.size >= Pt(24):
                    is_heading = True

                if is_heading:
                    lines.append(f"\n### {text}\n")
                elif para.level == 0:
                    lines.append(f"\n{text}\n")
                else:
                    indent = "  " * para.level
                    lines.append(f"{indent}- {text}")

        # Capture speaker notes
        notes_slide = slide.notes_slide
        if notes_slide:
            notes_text = notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append(f"\n> **Notes:** {notes_text}\n")

    result = "\n".join(lines)

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(result)
        print(f"Extracted {len(prs.slides)} slides → {output_path}")
    else:
        print(result)

    return result


if __name__ == "__main__":
    check_dependency()

    if len(sys.argv) < 2:
        print(f"Usage: python3 {sys.argv[0]} <input.pptx> [output.md]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(pptx_path):
        print(f"ERROR: File not found: {pptx_path}")
        sys.exit(1)

    extract(pptx_path, output_path)
